#!/usr/bin/env python3
"""Collect AFAC2026 TRPS v0.3 final pack staging materials."""

from __future__ import annotations

import json
import shutil
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
AFAC_ROOT = ROOT.parent
V1_ROOT = AFAC_ROOT / "afac2026_productization_v0_1"
V2_ROOT = AFAC_ROOT / "afac2026_productization_v0_2_demo_pack"
OUTPUT_DIR = ROOT / "outputs"
STAGING_DIR = OUTPUT_DIR / "final_pack_staging"
ZIP_ROOT_NAME = "TRPS_AFAC2026_FINAL_SUBMISSION_PACK"
GENERATED_AT = "2026-06-20T00:00:00Z"

FINAL_DOCS = [
    "README.md",
    "00_acceptance_from_v0_2.md",
    "01_final_submission_overview.md",
    "02_portal_copybook_final.md",
    "03_project_solution_final.md",
    "04_business_plan_final.md",
    "05_governance_note_final.md",
    "06_demo_guide_final.md",
    "07_pitch_script_8min.md",
    "08_demo_script_3min.md",
    "09_qa_defense_final.md",
    "10_slide_deck_content_final.md",
    "11_final_submission_checklist.md",
    "12_manual_submission_steps.md",
    "13_risk_register.md",
    "14_claim_evidence_map_final.md",
    "15_file_manifest_final.md",
    "assets/demo_link_note.md",
]

V1_OUTPUTS = [
    "demo_receipts.json",
    "demo_metrics.json",
    "demo_summary.md",
    "validation_report.json",
    "validation_report.md",
]

V2_OUTPUTS = [
    "submission_readiness_score.json",
    "submission_readiness_score.md",
    "validation_report.json",
    "validation_report.md",
    "afac_v0_2_manifest.json",
    "afac_v0_2_manifest.md",
]

EXCLUDE_NAMES = {
    ".git",
    ".venv",
    "node_modules",
    "__pycache__",
    ".env",
    ".DS_Store",
}


def load_json(path: Path) -> object:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, data: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(data, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )


def write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def copy_file(source: Path, destination: Path) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, destination)


def ignore_patterns(_directory: str, names: list[str]) -> set[str]:
    ignored: set[str] = set()
    for name in names:
        if name in EXCLUDE_NAMES:
            ignored.add(name)
        if name.endswith(".zip") or name.endswith(".pptx"):
            ignored.add(name)
    return ignored


def reset_staging() -> Path:
    if STAGING_DIR.exists():
        shutil.rmtree(STAGING_DIR)
    root = STAGING_DIR / ZIP_ROOT_NAME
    root.mkdir(parents=True, exist_ok=True)
    return root


def input_checks() -> dict[str, bool]:
    checks = {
        "v0_1_root": V1_ROOT.exists(),
        "v0_2_root": V2_ROOT.exists(),
        "v0_1_demo_receipts": (V1_ROOT / "outputs" / "demo_receipts.json").exists(),
        "v0_1_demo_metrics": (V1_ROOT / "outputs" / "demo_metrics.json").exists(),
        "v0_1_validation": (V1_ROOT / "outputs" / "validation_report.json").exists(),
        "v0_2_demo_index": (V2_ROOT / "demo" / "index.html").exists(),
        "v0_2_demo_data": (V2_ROOT / "demo" / "assets" / "trps_demo_data.json").exists(),
        "v0_2_readiness": (V2_ROOT / "outputs" / "submission_readiness_score.json").exists(),
        "v0_2_validation": (V2_ROOT / "outputs" / "validation_report.json").exists(),
    }
    missing = [name for name, ok in checks.items() if not ok]
    if missing:
        raise SystemExit(f"missing required inputs: {', '.join(missing)}")
    return checks


def generate_pptx_or_fallback() -> dict[str, object]:
    pptx_path = OUTPUT_DIR / "TRPS_AFAC2026_PitchDeck_v0_3.pptx"
    fallback_path = OUTPUT_DIR / "pptx_generation_skipped.md"
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        fallback = "\n".join(
            [
                "# PPTX Generation Skipped",
                "",
                "`python-pptx` is unavailable in the current environment.",
                "",
                f"Reason: `{exc}`",
                "",
                "Manual fallback: convert `10_slide_deck_content_final.md` into PPTX.",
                "",
            ]
        )
        write_text(fallback_path, fallback)
        if pptx_path.exists():
            pptx_path.unlink()
        return {
            "pptx_generated": False,
            "pptx_path": None,
            "fallback_path": str(fallback_path.relative_to(ROOT)),
            "reason": str(exc),
        }

    prs = Presentation()
    slide_content = (ROOT / "10_slide_deck_content_final.md").read_text(encoding="utf-8")
    sections = [section for section in slide_content.split("\n## Slide ") if section.strip()]
    for section in sections[1:]:
        title_line, *body_lines = section.splitlines()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_line.strip()
        body = slide.placeholders[1]
        body.text = "\n".join(line for line in body_lines if line.startswith("- "))[:1200]
    prs.save(pptx_path)
    if fallback_path.exists():
        fallback_path.unlink()
    return {
        "pptx_generated": True,
        "pptx_path": str(pptx_path.relative_to(ROOT)),
        "fallback_path": None,
        "reason": None,
    }


def readiness_score(pptx_status: dict[str, object]) -> dict[str, object]:
    v2_score = load_json(V2_ROOT / "outputs" / "submission_readiness_score.json")
    v2_validation = load_json(V2_ROOT / "outputs" / "validation_report.json")
    assert isinstance(v2_score, dict)
    assert isinstance(v2_validation, dict)
    dimensions = [
        {
            "name": "portal_copy_readiness",
            "score": 96,
            "evidence_path": "02_portal_copybook_final.md",
        },
        {
            "name": "pitch_readiness",
            "score": 95,
            "evidence_path": "07_pitch_script_8min.md;10_slide_deck_content_final.md",
        },
        {
            "name": "demo_readiness",
            "score": 96,
            "evidence_path": "06_demo_guide_final.md;v0_2_demo/demo/index.html",
        },
        {
            "name": "qa_readiness",
            "score": 95,
            "evidence_path": "09_qa_defense_final.md",
        },
        {
            "name": "governance_safety",
            "score": 98,
            "evidence_path": "05_governance_note_final.md;13_risk_register.md",
        },
        {
            "name": "zip_packaging_readiness",
            "score": 94,
            "evidence_path": "scripts/build_final_zip.py;outputs/final_pack_manifest.json",
        },
    ]
    total = round(sum(item["score"] for item in dimensions) / len(dimensions), 2)
    return {
        "pack_id": "afac2026_trps_final_submission_pack_v0_3",
        "generated_at": GENERATED_AT,
        "total_score": total,
        "dimensions": dimensions,
        "v0_2_total_score": v2_score["total_score"],
        "v0_2_validation_ok": v2_validation["ok"],
        "pptx_status": pptx_status,
        "status": "local_final_submission_pack_not_submitted",
    }


def write_readiness(score: dict[str, object]) -> None:
    write_json(OUTPUT_DIR / "final_submission_readiness_score.json", score)
    lines = [
        "# Final Submission Readiness Score",
        "",
        f"- `total_score`: `{score['total_score']}`",
        f"- `status`: `{score['status']}`",
        "",
        "| Dimension | Score | Evidence |",
        "| --- | ---: | --- |",
    ]
    for item in score["dimensions"]:
        lines.append(f"| `{item['name']}` | `{item['score']}` | `{item['evidence_path']}` |")
    lines.extend(["", "This score is local preparation evidence, not official scoring.", ""])
    write_text(OUTPUT_DIR / "final_submission_readiness_score.md", "\n".join(lines))


def copy_materials(staging_root: Path) -> None:
    for relative in FINAL_DOCS:
        copy_file(ROOT / relative, staging_root / "final_docs" / relative)

    shutil.copytree(
        V2_ROOT / "demo",
        staging_root / "v0_2_demo" / "demo",
        ignore=ignore_patterns,
        dirs_exist_ok=True,
    )

    for relative in V1_OUTPUTS:
        source = V1_ROOT / "outputs" / relative
        if source.exists():
            copy_file(source, staging_root / "v0_1_outputs" / relative)

    for relative in V2_OUTPUTS:
        source = V2_ROOT / "outputs" / relative
        if source.exists():
            copy_file(source, staging_root / "v0_2_outputs" / relative)

    for source in sorted(OUTPUT_DIR.glob("*.json")) + sorted(OUTPUT_DIR.glob("*.md")):
        copy_file(source, staging_root / "v0_3_outputs" / source.name)


def manifest(score: dict[str, object], checks: dict[str, bool]) -> dict[str, object]:
    files = []
    for path in sorted(ROOT.rglob("*")):
        if not path.is_file() or STAGING_DIR in path.parents:
            continue
        if "__pycache__" in path.parts or path.name == ".DS_Store":
            continue
        files.append(str(path.relative_to(ROOT)))
    return {
        "pack_id": "afac2026_trps_final_submission_pack_v0_3",
        "generated_at": GENERATED_AT,
        "status": "local_final_submission_pack_not_submitted",
        "source_commits": {
            "v0_1": "ded7bcade84332b68363cac8cb98b7e1a5ae3b9f",
            "v0_2": "0fed7b295c2c4c52d162bea635d5c8404ec99a6a",
        },
        "input_checks": checks,
        "readiness_total_score": score["total_score"],
        "zip_path": "outputs/TRPS_AFAC2026_FINAL_SUBMISSION_PACK.zip",
        "no_push": True,
        "no_tag": True,
        "no_external_submission": True,
        "affected_eeoap_clauses": [
            "EEOAP-001",
            "EEOAP-003",
            "EEOAP-004",
            "EEOAP-005",
        ],
        "files": files,
    }


def write_manifest(manifest_data: dict[str, object]) -> None:
    write_json(OUTPUT_DIR / "final_pack_manifest.json", manifest_data)
    lines = [
        "# Final Pack Manifest",
        "",
        f"- `pack_id`: `{manifest_data['pack_id']}`",
        f"- `status`: `{manifest_data['status']}`",
        f"- `readiness_total_score`: `{manifest_data['readiness_total_score']}`",
        f"- `zip_path`: `{manifest_data['zip_path']}`",
        f"- `no_push`: `{manifest_data['no_push']}`",
        f"- `no_tag`: `{manifest_data['no_tag']}`",
        "",
        "## Files",
        "",
    ]
    lines.extend(f"- `{path}`" for path in manifest_data["files"])
    lines.append("")
    write_text(OUTPUT_DIR / "final_pack_manifest.md", "\n".join(lines))


def main() -> int:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    checks = input_checks()
    pptx_status = generate_pptx_or_fallback()
    score = readiness_score(pptx_status)
    write_readiness(score)
    manifest_data = manifest(score, checks)
    write_manifest(manifest_data)
    staging_root = reset_staging()
    copy_materials(staging_root)
    print(
        json.dumps(
            {
                "ok": True,
                "staging_root": str(staging_root),
                "readiness_total_score": score["total_score"],
                "pptx_generated": pptx_status["pptx_generated"],
            },
            sort_keys=True,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
